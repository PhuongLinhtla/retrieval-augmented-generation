from pathlib import Path
from collections import Counter
import re

from pypdf import PdfReader

from .schemas import TextSpan
from .vietnamese_nlp import (
    estimate_token_count,
    extract_keywords,
    fix_common_ocr_errors,
    normalize_whitespace,
    split_sentences_vietnamese,
    split_text_units,
)


SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".markdown", ".log", ".pptx"}


def clean_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = normalize_whitespace(text)
    text = fix_common_ocr_errors(text)
    return text.strip()


def load_text_spans(file_path: Path) -> list[TextSpan]:
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file format: {suffix}")

    if suffix == ".pdf":
        return _load_pdf(file_path)
    if suffix == ".pptx":
        return _load_pptx(file_path)
    return _load_plain_text(file_path)


def chunk_spans(
    spans: list[TextSpan],
    chunk_size: int,
    chunk_overlap: int,
) -> list[tuple[int, int, str]]:
    chunks: list[tuple[int, int, str]] = []
    chunk_index = 0

    for span in spans:
        units = _build_units(span.text, max(220, chunk_size // 2))
        for chunk_text in _pack_units_to_chunks(units, chunk_size, chunk_overlap):
            chunks.append((span.page_number, chunk_index, chunk_text))
            chunk_index += 1

    return chunks


def build_parent_child_chunks(
    spans: list[TextSpan],
    child_chunk_size: int,
    child_chunk_overlap: int,
    parent_chunk_size: int,
) -> tuple[list[tuple[int, int, str, str, str, str]], list[tuple[int, int, int, str]]]:
    return build_parent_child_chunks_with_tokens(
        spans=spans,
        child_chunk_target_tokens=max(220, child_chunk_size // 2),
        child_chunk_min_tokens=max(120, child_chunk_size // 3),
        child_chunk_max_tokens=max(260, child_chunk_size),
        child_chunk_overlap_sentences=max(1, child_chunk_overlap // 60),
        parent_chunk_target_tokens=max(700, parent_chunk_size // 2),
        parent_chunk_min_tokens=max(420, parent_chunk_size // 3),
        parent_chunk_max_tokens=max(900, parent_chunk_size),
        parent_child_group_min=2,
        parent_child_group_max=4,
    )


def build_parent_child_chunks_with_tokens(
    spans: list[TextSpan],
    *,
    child_chunk_target_tokens: int,
    child_chunk_min_tokens: int,
    child_chunk_max_tokens: int,
    child_chunk_overlap_sentences: int,
    parent_chunk_target_tokens: int,
    parent_chunk_min_tokens: int,
    parent_chunk_max_tokens: int,
    parent_child_group_min: int,
    parent_child_group_max: int,
    chapter_hint: str = "",
) -> tuple[list[tuple[int, int, str, str, str, str]], list[tuple[int, int, int, str]]]:
    parent_rows: list[tuple[int, int, str, str, str, str]] = []
    child_rows: list[tuple[int, int, int, str]] = []

    parent_index = 0
    child_index = 0

    for span in spans:
        section_rows = _split_sections(span.text)
        for section_title, section_text in section_rows:
            blocks = _build_semantic_blocks(
                section_text,
                sentence_pack_max=max(24, child_chunk_target_tokens // 4),
            )
            if not blocks:
                continue

            child_chunks = _pack_blocks_to_child_chunks(
                blocks=blocks,
                target_tokens=child_chunk_target_tokens,
                min_tokens=child_chunk_min_tokens,
                max_tokens=child_chunk_max_tokens,
                overlap_sentences=child_chunk_overlap_sentences,
            )
            if not child_chunks:
                continue

            parent_groups = _group_children_for_parent(
                child_chunks=child_chunks,
                target_tokens=parent_chunk_target_tokens,
                min_tokens=parent_chunk_min_tokens,
                max_tokens=parent_chunk_max_tokens,
                min_group_size=max(1, parent_child_group_min),
                max_group_size=max(1, parent_child_group_max),
            )

            for group in parent_groups:
                parent_text = clean_text("\n\n".join(group))
                if not parent_text:
                    continue

                current_parent_index = parent_index
                section_path = _build_section_path(chapter_hint=chapter_hint, section_title=section_title)
                keywords = ",".join(extract_keywords(parent_text))
                parent_rows.append(
                    (
                        span.page_number,
                        current_parent_index,
                        _derive_parent_title(parent_text, section_hint=section_title),
                        section_path,
                        keywords,
                        parent_text,
                    )
                )
                parent_index += 1

                for child_text in group:
                    cleaned_child = clean_text(child_text)
                    if not cleaned_child:
                        continue
                    child_rows.append(
                        (
                            span.page_number,
                            child_index,
                            current_parent_index,
                            cleaned_child,
                        )
                    )
                    child_index += 1

    return parent_rows, child_rows


def _split_sections(text: str) -> list[tuple[str, str]]:
    cleaned = clean_text(text)
    if not cleaned:
        return []

    lines = [line.rstrip() for line in cleaned.splitlines()]
    sections: list[tuple[str, str]] = []
    current_title = "Section"
    current_lines: list[str] = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            if current_lines:
                current_lines.append("")
            continue

        if _is_heading_line(stripped):
            if current_lines:
                section_text = clean_text("\n".join(current_lines))
                if section_text:
                    sections.append((current_title, section_text))
                current_lines.clear()

            current_title = re.sub(r"^#{1,6}\s*", "", stripped).strip()
            continue

        current_lines.append(stripped)

    if current_lines:
        section_text = clean_text("\n".join(current_lines))
        if section_text:
            sections.append((current_title, section_text))

    if not sections:
        return [("Section", cleaned)]
    return sections


def _is_heading_line(line: str) -> bool:
    return bool(
        re.match(
            r"^(#{1,6}\s+|chuong\s+\d+|chapter\s+\d+|muc\s+\d+|section\s+\d+|\d+(\.\d+){0,3}\s+)",
            line,
            re.IGNORECASE,
        )
    )


def _build_semantic_blocks(section_text: str, sentence_pack_max: int) -> list[str]:
    blocks: list[str] = []
    bullet_buffer: list[str] = []

    def flush_bullets() -> None:
        if not bullet_buffer:
            return
        merged = clean_text("\n".join(bullet_buffer))
        if merged:
            blocks.append(merged)
        bullet_buffer.clear()

    for unit in split_text_units(section_text):
        text = clean_text(unit.text)
        if not text:
            continue

        if unit.kind == "bullet":
            bullet_buffer.append(text)
            continue

        flush_bullets()

        if unit.kind == "formula":
            blocks.append(text)
            continue

        if unit.kind == "paragraph":
            sentences = split_sentences_vietnamese(text)
            blocks.extend(_pack_sentences(sentences, sentence_pack_max))
            continue

        blocks.append(text)

    flush_bullets()
    return _merge_question_answer_blocks(blocks)


def _pack_sentences(sentences: list[str], max_tokens: int) -> list[str]:
    if not sentences:
        return []

    packed: list[str] = []
    current = ""

    for sentence in sentences:
        sentence = clean_text(sentence)
        if not sentence:
            continue

        if not current:
            current = sentence
            continue

        candidate = f"{current} {sentence}".strip()
        if estimate_token_count(candidate) <= max_tokens:
            current = candidate
        else:
            packed.append(current)
            current = sentence

    if current:
        packed.append(current)
    return packed


def _merge_question_answer_blocks(blocks: list[str]) -> list[str]:
    if not blocks:
        return []

    merged: list[str] = []
    idx = 0

    while idx < len(blocks):
        current = blocks[idx]
        if (
            idx + 1 < len(blocks)
            and _is_question_block(current)
            and _is_answer_block(blocks[idx + 1])
        ):
            merged.append(clean_text(f"{current}\n{blocks[idx + 1]}"))
            idx += 2
            continue

        merged.append(current)
        idx += 1

    return merged


def _is_question_block(text: str) -> bool:
    lowered = text.casefold()
    return lowered.startswith("cau hoi") or lowered.startswith("question")


def _is_answer_block(text: str) -> bool:
    lowered = text.casefold()
    return lowered.startswith("loi giai") or lowered.startswith("dap an") or lowered.startswith("answer")


def _pack_blocks_to_child_chunks(
    blocks: list[str],
    target_tokens: int,
    min_tokens: int,
    max_tokens: int,
    overlap_sentences: int,
) -> list[str]:
    if not blocks:
        return []

    chunks: list[str] = []
    buffer: list[str] = []

    def flush_buffer() -> None:
        if not buffer:
            return
        chunk_text = clean_text("\n".join(buffer))
        if chunk_text:
            chunks.append(chunk_text)

    for block in blocks:
        block = clean_text(block)
        if not block:
            continue

        block_tokens = estimate_token_count(block)
        if block_tokens > max_tokens and not _is_protected_block(block):
            split_blocks = _split_block_by_tokens(block, max_tokens)
        else:
            split_blocks = [block]

        for piece in split_blocks:
            piece_tokens = estimate_token_count(piece)
            current_tokens = estimate_token_count("\n".join(buffer)) if buffer else 0

            if buffer and current_tokens + piece_tokens > max_tokens:
                flush_buffer()
                previous = chunks[-1] if chunks else ""
                overlap_text = _take_sentence_overlap(previous, overlap_sentences)
                buffer.clear()
                if overlap_text:
                    buffer.append(overlap_text)

            buffer.append(piece)
            current_tokens = estimate_token_count("\n".join(buffer))
            if current_tokens >= target_tokens and current_tokens >= min_tokens:
                flush_buffer()
                previous = chunks[-1] if chunks else ""
                overlap_text = _take_sentence_overlap(previous, overlap_sentences)
                buffer.clear()
                if overlap_text:
                    buffer.append(overlap_text)

    flush_buffer()

    if len(chunks) >= 2 and estimate_token_count(chunks[-1]) < max(20, min_tokens // 2):
        tail = chunks.pop()
        chunks[-1] = clean_text(f"{chunks[-1]}\n{tail}")

    return chunks


def _is_protected_block(block: str) -> bool:
    if "\n- " in block or "\n* " in block:
        return True
    if _is_question_block(block) or _is_answer_block(block):
        return True
    return bool(re.search(r"[=<>^]{1,}", block))


def _split_block_by_tokens(text: str, max_tokens: int) -> list[str]:
    sentences = split_sentences_vietnamese(text)
    if len(sentences) <= 1:
        return _split_by_chars(text, max(120, max_tokens * 4))
    return _pack_sentences(sentences, max_tokens)


def _take_sentence_overlap(text: str, overlap_sentences: int) -> str:
    if overlap_sentences <= 0:
        return ""
    sentences = split_sentences_vietnamese(text)
    if len(sentences) <= 1:
        return ""
    return clean_text(" ".join(sentences[-overlap_sentences:]))


def _group_children_for_parent(
    child_chunks: list[str],
    target_tokens: int,
    min_tokens: int,
    max_tokens: int,
    min_group_size: int,
    max_group_size: int,
) -> list[list[str]]:
    if not child_chunks:
        return []

    groups: list[list[str]] = []
    idx = 0

    while idx < len(child_chunks):
        remaining = len(child_chunks) - idx
        best_size = min(max_group_size, remaining)
        best_score = float("inf")

        for size in range(max(1, min_group_size), min(max_group_size, remaining) + 1):
            candidate = child_chunks[idx : idx + size]
            token_count = estimate_token_count("\n\n".join(candidate))
            if token_count > max_tokens and size > 1:
                continue

            penalty = abs(token_count - target_tokens)
            if token_count < min_tokens:
                penalty += min_tokens - token_count

            if penalty < best_score:
                best_score = penalty
                best_size = size

        group = list(child_chunks[idx : idx + best_size])
        if not group:
            break

        if groups and remaining < min_group_size:
            groups[-1].extend(group)
        else:
            groups.append(group)
        idx += best_size

    return groups


def _build_section_path(chapter_hint: str, section_title: str) -> str:
    chapter_text = clean_text(chapter_hint)
    section_text = clean_text(section_title)

    if chapter_text and section_text and chapter_text.casefold() != section_text.casefold():
        return f"{chapter_text} > {section_text}"
    if section_text:
        return section_text
    if chapter_text:
        return chapter_text
    return ""


def _derive_parent_title(text: str, section_hint: str = "", max_len: int = 120) -> str:
    section_hint = clean_text(section_hint)
    if section_hint and section_hint.casefold() not in {"section", "chuong", "chapter"}:
        title = section_hint
    else:
        first_line = clean_text(text.split("\n", 1)[0])
        if _is_heading_line(first_line):
            title = re.sub(r"^#{1,6}\s*", "", first_line).strip()
        else:
            sentence_candidates = split_sentences_vietnamese(first_line)
            title = sentence_candidates[0] if sentence_candidates else first_line

    if not title:
        return "Section"
    if len(title) <= max_len:
        return title
    return f"{title[: max_len - 3].rstrip()}..."


def _load_pdf(file_path: Path) -> list[TextSpan]:
    reader = PdfReader(str(file_path))
    raw_pages: list[str] = []

    for page in reader.pages:
        raw_pages.append(clean_text(page.extract_text() or ""))

    cleaned_pages = _remove_repeated_page_noise(raw_pages)
    spans: list[TextSpan] = []
    for index, page_text in enumerate(cleaned_pages, start=1):
        if page_text:
            spans.append(TextSpan(page_number=index, text=page_text))

    return spans


def _remove_repeated_page_noise(page_texts: list[str]) -> list[str]:
    if len(page_texts) < 3:
        return page_texts

    top_counter: Counter[str] = Counter()
    bottom_counter: Counter[str] = Counter()

    parsed_lines: list[list[str]] = []
    for text in page_texts:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        parsed_lines.append(lines)
        if lines:
            top_counter[lines[0]] += 1
            bottom_counter[lines[-1]] += 1

    threshold = max(2, int(len(page_texts) * 0.6))
    repeated_top = {
        line
        for line, count in top_counter.items()
        if count >= threshold and len(line) <= 120
    }
    repeated_bottom = {
        line
        for line, count in bottom_counter.items()
        if count >= threshold and len(line) <= 120
    }

    cleaned: list[str] = []
    for lines in parsed_lines:
        filtered = list(lines)
        if filtered and filtered[0] in repeated_top:
            filtered = filtered[1:]
        if filtered and filtered[-1] in repeated_bottom:
            filtered = filtered[:-1]
        cleaned.append(clean_text("\n".join(filtered)))

    return cleaned


def _load_plain_text(file_path: Path) -> list[TextSpan]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    text = clean_text(text)
    if not text:
        return []
    return [TextSpan(page_number=1, text=text)]


def _load_pptx(file_path: Path) -> list[TextSpan]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for .pptx ingestion") from exc

    presentation = Presentation(str(file_path))
    spans: list[TextSpan] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        fragments: list[str] = []

        for shape in slide.shapes:
            text = ""
            if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None):
                text = str(shape.text_frame.text or "")
            elif hasattr(shape, "text"):
                text = str(getattr(shape, "text") or "")

            text = clean_text(text)
            if text:
                fragments.append(text)

        content = clean_text("\n".join(fragments))
        if content:
            spans.append(TextSpan(page_number=slide_index, text=content))

    return spans


def _build_units(text: str, max_unit_chars: int) -> list[str]:
    paragraphs = [clean_text(par) for par in re.split(r"\n\n+", text)]
    paragraphs = [par for par in paragraphs if par]

    units: list[str] = []
    for paragraph in paragraphs:
        if len(paragraph) <= max_unit_chars:
            units.append(paragraph)
            continue
        units.extend(_split_long_paragraph(paragraph, max_unit_chars))

    return units


def _build_parent_chunks(text: str, max_parent_chars: int) -> list[str]:
    cleaned = clean_text(text)
    if not cleaned:
        return []

    paragraphs = [clean_text(par) for par in re.split(r"\n\n+", cleaned)]
    paragraphs = [par for par in paragraphs if par]
    if not paragraphs:
        return _split_by_chars(cleaned, max_parent_chars)

    grouped: list[str] = []
    current = ""

    for paragraph in paragraphs:
        if not current:
            current = paragraph
            continue

        if len(current) + 2 + len(paragraph) <= max_parent_chars:
            current = f"{current}\n\n{paragraph}"
        else:
            grouped.append(current.strip())
            current = paragraph

    if current:
        grouped.append(current.strip())

    finalized: list[str] = []
    for chunk in grouped:
        if len(chunk) <= max_parent_chars:
            finalized.append(chunk)
            continue
        finalized.extend(_split_long_paragraph(chunk, max_parent_chars))

    return finalized


def _derive_parent_title(text: str, max_len: int = 120) -> str:
    if not text:
        return "Section"

    first_line = clean_text(text.split("\n", 1)[0])
    candidates = [candidate.strip() for candidate in re.split(r"(?<=[.!?])\s+", first_line)]
    title = next((candidate for candidate in candidates if candidate), first_line)

    if not title:
        return "Section"
    if len(title) <= max_len:
        return title
    return f"{title[: max_len - 3].rstrip()}..."


def _split_long_paragraph(paragraph: str, max_unit_chars: int) -> list[str]:
    sentences = re.split(r"(?<=[.!?])\s+", paragraph)
    sentences = [sentence.strip() for sentence in sentences if sentence.strip()]

    if not sentences:
        return _split_by_chars(paragraph, max_unit_chars)

    out: list[str] = []
    current = ""
    for sentence in sentences:
        if not current:
            current = sentence
            continue

        if len(current) + 1 + len(sentence) <= max_unit_chars:
            current = f"{current} {sentence}"
        else:
            out.append(current)
            current = sentence

    if current:
        out.append(current)

    final_units: list[str] = []
    for unit in out:
        if len(unit) <= max_unit_chars:
            final_units.append(unit)
        else:
            final_units.extend(_split_by_chars(unit, max_unit_chars))

    return final_units


def _split_by_chars(text: str, max_unit_chars: int) -> list[str]:
    slices: list[str] = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(text_length, start + max_unit_chars)
        if end < text_length:
            split_pos = text.rfind(" ", start, end)
            if split_pos > start + max_unit_chars // 2:
                end = split_pos

        part = text[start:end].strip()
        if part:
            slices.append(part)

        if end == text_length:
            break
        start = end

    return slices


def _pack_units_to_chunks(
    units: list[str],
    chunk_size: int,
    chunk_overlap: int,
) -> list[str]:
    if not units:
        return []

    chunks: list[str] = []
    buffer: list[str] = []
    buffer_length = 0

    for unit in units:
        unit = unit.strip()
        if not unit:
            continue

        extra_length = len(unit) + (1 if buffer else 0)
        if buffer and buffer_length + extra_length > chunk_size:
            chunk_text = " ".join(buffer).strip()
            if chunk_text:
                chunks.append(chunk_text)

            overlap_text = _take_overlap(chunk_text, chunk_overlap)
            if overlap_text:
                buffer = [overlap_text, unit]
            else:
                buffer = [unit]
            buffer_length = sum(len(part) for part in buffer) + max(0, len(buffer) - 1)
        else:
            buffer.append(unit)
            buffer_length += extra_length

    if buffer:
        chunk_text = " ".join(buffer).strip()
        if chunk_text:
            chunks.append(chunk_text)

    return chunks


def _take_overlap(text: str, overlap_chars: int) -> str:
    if overlap_chars <= 0 or len(text) <= overlap_chars:
        return ""

    tail = text[-overlap_chars:].strip()
    first_space = tail.find(" ")
    if first_space > 0:
        tail = tail[first_space + 1 :]
    return tail.strip()
